from __future__ import annotations

import asyncio
import os
import tempfile
import time
import uuid
import zipfile
from pathlib import Path
from typing import Any

import anyio
from dotenv import load_dotenv, set_key
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse, HTMLResponse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from magical_layers.agentic import run_agentic_selection
from magical_layers.evaluation import compare_render, render_pptx_first_slide, render_pptx_slides
from magical_layers.llm import DEFAULT_FREE_BRAIN_MODELS
from magical_layers.openrouter_judge import judge_model_chain
from magical_layers.orchestrator import PipelineOptions, image_to_layers
from magical_layers.presets import segment_preset
from magical_layers.pptx_writer import write_pptx, write_pptx_deck


ENV_PATH = Path(".env")
load_dotenv(ENV_PATH)

app = FastAPI(title="Magical Layers Local")

JOBS: dict[str, dict[str, Any]] = {}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
ZIP_MEDIA_TYPE = "application/zip"


@app.get("/", response_class=HTMLResponse)
def index() -> str:
    openrouter_available = "true" if os.getenv("OPENROUTER_API_KEY") else "false"
    return f"""
<!doctype html>
<html lang="es">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Magical Layers Local</title>
  <style>
    :root {{
      color-scheme: light;
      --bg: #f4f7f6;
      --ink: #17201d;
      --muted: #65726d;
      --line: #d9e1de;
      --panel: #ffffff;
      --panel-soft: #eef4f1;
      --accent: #0d8f73;
      --accent-strong: #08745d;
      --warm: #c77620;
      --danger: #a33a32;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      min-height: 100vh;
      font-family: Inter, ui-sans-serif, system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
      background: var(--bg);
      color: var(--ink);
      letter-spacing: 0;
    }}
    main {{
      width: min(1180px, calc(100vw - 32px));
      margin: 0 auto;
      padding: 28px 0 32px;
    }}
    header {{
      display: flex;
      align-items: end;
      justify-content: space-between;
      gap: 20px;
      margin-bottom: 18px;
    }}
    h1 {{
      margin: 0;
      font-size: clamp(30px, 5vw, 58px);
      line-height: .92;
      font-weight: 820;
    }}
    .subtitle {{
      margin: 8px 0 0;
      max-width: 680px;
      color: var(--muted);
      font-size: 15px;
    }}
    .quality-badge {{
      min-width: 198px;
      border: 1px solid var(--line);
      background: var(--panel);
      border-radius: 8px;
      padding: 12px 14px;
    }}
    .quality-badge strong {{
      display: block;
      font-size: 30px;
      line-height: 1;
    }}
    .quality-badge span {{
      display: block;
      margin-top: 5px;
      color: var(--muted);
      font-size: 12px;
    }}
    .workspace {{
      display: grid;
      grid-template-columns: minmax(0, 1.15fr) minmax(330px, .85fr);
      gap: 16px;
      align-items: start;
    }}
    section {{
      border: 1px solid var(--line);
      background: var(--panel);
      border-radius: 8px;
    }}
    .tool {{
      padding: 18px;
    }}
    .tabs {{
      display: grid;
      grid-template-columns: repeat(2, minmax(0, 1fr));
      gap: 6px;
      margin-bottom: 14px;
      padding: 5px;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: var(--panel-soft);
    }}
    .tab-button {{
      min-height: 38px;
      background: transparent;
      color: var(--muted);
      box-shadow: none;
    }}
    .tab-button:hover {{
      background: #f7faf9;
      color: var(--ink);
    }}
    .tab-button.active {{
      background: var(--panel);
      color: var(--ink);
      box-shadow: inset 0 0 0 1px #cfdad6;
    }}
    .tab-panel[hidden] {{ display: none; }}
    .dropzone {{
      position: relative;
      display: grid;
      place-items: center;
      min-height: 236px;
      border: 1px dashed #93a29d;
      border-radius: 8px;
      background: #f9fbfa;
      overflow: hidden;
    }}
    .dropzone input {{
      position: absolute;
      inset: 0;
      opacity: 0;
      cursor: pointer;
    }}
    .drop-inner {{
      width: min(430px, calc(100% - 24px));
      text-align: center;
      pointer-events: none;
    }}
    .drop-title {{
      margin: 0;
      font-size: 22px;
      font-weight: 760;
    }}
    .drop-meta {{
      margin: 8px 0 0;
      color: var(--muted);
      font-size: 14px;
    }}
    .filename {{
      display: inline-flex;
      max-width: 100%;
      margin-top: 15px;
      padding: 7px 10px;
      border: 1px solid var(--line);
      border-radius: 999px;
      color: var(--ink);
      background: var(--panel);
      font-size: 13px;
      white-space: nowrap;
      overflow: hidden;
      text-overflow: ellipsis;
    }}
    .controls {{
      display: grid;
      grid-template-columns: repeat(2, minmax(0, 1fr));
      gap: 10px;
      margin-top: 14px;
    }}
    .check {{
      min-height: 72px;
      display: grid;
      grid-template-columns: 22px 1fr;
      gap: 10px;
      align-items: start;
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 12px;
      background: #fbfcfc;
    }}
    .check input {{
      width: 18px;
      height: 18px;
      margin: 1px 0 0;
      accent-color: var(--accent);
    }}
    .check label {{
      display: block;
      font-weight: 720;
      line-height: 1.15;
    }}
    .check small {{
      display: block;
      margin-top: 4px;
      color: var(--muted);
      font-size: 12px;
      line-height: 1.25;
    }}
    .check select {{
      width: 100%;
      min-height: 34px;
      margin-top: 8px;
      border: 1px solid var(--line);
      border-radius: 6px;
      padding: 0 8px;
      background: var(--panel);
      color: var(--ink);
      font: inherit;
      font-size: 12px;
    }}
    .segmented {{
      display: grid;
      grid-template-columns: repeat(3, minmax(0, 1fr));
      gap: 6px;
      margin-top: 14px;
      padding: 5px;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: var(--panel-soft);
    }}
    .segmented input {{
      position: absolute;
      opacity: 0;
      pointer-events: none;
    }}
    .segmented label {{
      display: grid;
      place-items: center;
      min-height: 38px;
      border-radius: 6px;
      color: var(--muted);
      font-weight: 720;
      cursor: pointer;
      white-space: nowrap;
    }}
    .segmented input:checked + label {{
      color: var(--ink);
      background: var(--panel);
      box-shadow: inset 0 0 0 1px #cfdad6;
    }}
    .output-mode {{
      grid-template-columns: repeat(2, minmax(0, 1fr));
    }}
    .actions {{
      display: flex;
      align-items: center;
      gap: 10px;
      margin-top: 14px;
    }}
    button,
    .download,
    .small-button {{
      min-height: 44px;
      border: 0;
      border-radius: 8px;
      padding: 0 15px;
      font: inherit;
      font-weight: 780;
      text-decoration: none;
      display: inline-grid;
      place-items: center;
      cursor: pointer;
    }}
    button {{
      color: white;
      background: var(--accent);
    }}
    button:hover {{ background: var(--accent-strong); }}
    button:disabled {{
      cursor: progress;
      background: #8aa29a;
    }}
    .download {{
      color: var(--ink);
      background: #f7e4ca;
      border: 1px solid #ecc48e;
    }}
    .download[hidden] {{ display: none; }}
    .key-row {{
      display: grid;
      grid-template-columns: 1fr auto;
      gap: 7px;
      margin-top: 8px;
    }}
    .key-row[hidden] {{ display: none; }}
    .key-row input {{
      width: 100%;
      min-width: 0;
      height: 36px;
      border: 1px solid var(--line);
      border-radius: 7px;
      padding: 0 9px;
      font: inherit;
      font-size: 12px;
    }}
    .small-button {{
      min-height: 36px;
      padding: 0 10px;
      font-size: 12px;
      background: #15384c;
      color: white;
    }}
    .models-panel {{
      display: grid;
      gap: 12px;
    }}
    .models-panel label {{
      display: block;
      font-weight: 760;
    }}
    .models-panel textarea {{
      width: 100%;
      min-height: 190px;
      resize: vertical;
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 10px 11px;
      font: 13px ui-monospace, SFMono-Regular, Consolas, "Liberation Mono", monospace;
      color: var(--ink);
      background: #fbfcfc;
    }}
    .models-actions {{
      display: flex;
      align-items: center;
      gap: 10px;
      flex-wrap: wrap;
    }}
    .model-status {{
      color: var(--muted);
      font-size: 13px;
    }}
    .status {{
      padding: 18px;
    }}
    .meter-head {{
      display: flex;
      justify-content: space-between;
      gap: 12px;
      margin-bottom: 10px;
      color: var(--muted);
      font-size: 13px;
      font-weight: 680;
    }}
    .meter {{
      width: 100%;
      height: 12px;
      overflow: hidden;
      border-radius: 999px;
      background: #e2ebe7;
    }}
    .meter div {{
      width: 0%;
      height: 100%;
      border-radius: inherit;
      background: linear-gradient(90deg, var(--accent), #4aa686, var(--warm));
      transition: width .35s ease;
    }}
    .stage {{
      min-height: 24px;
      margin: 14px 0 0;
      font-size: 18px;
      font-weight: 760;
    }}
    .message {{
      min-height: 20px;
      margin: 5px 0 0;
      color: var(--muted);
      font-size: 13px;
    }}
    .steps {{
      display: grid;
      gap: 8px;
      margin-top: 16px;
    }}
    .step {{
      display: grid;
      grid-template-columns: 24px 1fr;
      gap: 9px;
      align-items: center;
      color: var(--muted);
      font-size: 13px;
    }}
    .dot {{
      width: 14px;
      height: 14px;
      border-radius: 50%;
      border: 2px solid #b6c4bf;
      justify-self: center;
      background: white;
    }}
    .step.done {{ color: var(--ink); }}
    .step.done .dot {{
      border-color: var(--accent);
      background: var(--accent);
    }}
    .metrics {{
      display: grid;
      grid-template-columns: repeat(3, minmax(0, 1fr));
      gap: 8px;
      margin-top: 16px;
    }}
    .metric {{
      min-height: 72px;
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 10px;
      background: #fbfcfc;
    }}
    .metric span {{
      display: block;
      color: var(--muted);
      font-size: 12px;
    }}
    .metric strong {{
      display: block;
      margin-top: 6px;
      font-size: 22px;
      line-height: 1;
    }}
    .result-line {{
      margin-top: 14px;
      color: var(--muted);
      font-size: 13px;
    }}
    .error {{
      color: var(--danger);
    }}
    footer {{
      display: grid;
      grid-template-columns: repeat(4, minmax(0, 1fr));
      gap: 10px;
      margin-top: 16px;
    }}
    .stat {{
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 11px 12px;
      background: var(--panel);
    }}
    .stat span {{
      display: block;
      color: var(--muted);
      font-size: 12px;
    }}
    .stat strong {{
      display: block;
      margin-top: 4px;
      font-size: 19px;
    }}
    @media (max-width: 860px) {{
      main {{ width: min(100% - 20px, 640px); padding-top: 18px; }}
      header, .workspace {{ display: grid; }}
      .quality-badge {{ min-width: 0; }}
      .workspace {{ grid-template-columns: 1fr; }}
      .controls, footer {{ grid-template-columns: 1fr; }}
      .metrics {{ grid-template-columns: 1fr; }}
      .segmented {{ grid-template-columns: 1fr; }}
      .actions {{ display: grid; }}
      button, .download {{ width: 100%; }}
    }}
  </style>
</head>
<body>
  <main>
    <header>
      <div>
        <h1>Magical Layers</h1>
        <p class="subtitle">Convierte una o varias imágenes planas en un PPTX con capas movibles, rápido y medible.</p>
      </div>
      <div class="quality-badge">
        <strong>A-/B+</strong>
        <span>24/24 procesadas · MAE 4.106 · Edge F1 0.864</span>
      </div>
    </header>

    <div class="workspace">
      <section class="tool" aria-label="Generador PPTX">
        <div class="tabs" role="tablist" aria-label="Paneles">
          <button class="tab-button active" type="button" data-tab="generate">Generar</button>
          <button class="tab-button" type="button" data-tab="models">Modelos</button>
        </div>

        <div class="tab-panel" id="panel-generate">
        <form id="job-form">
          <div class="dropzone" id="dropzone">
            <input id="file" name="files" type="file" accept="image/png,image/jpeg,image/webp" multiple required>
            <div class="drop-inner">
              <p class="drop-title">Suelta imágenes o selecciona archivos</p>
              <p class="drop-meta">PNG, JPG o WEBP · una slide por imagen</p>
              <span class="filename" id="filename">Ningún archivo seleccionado</span>
            </div>
          </div>

          <div class="controls">
            <div class="check">
              <input id="quality-check" name="quality_check" type="checkbox" value="true" checked>
              <div>
                <label for="quality-check">Medir calidad</label>
                <small>Renderiza el PPTX y compara contra la imagen.</small>
              </div>
            </div>
            <div class="check">
              <input id="ocr" name="ocr" type="checkbox" value="true">
              <div>
                <label for="ocr">Texto OCR editable</label>
                <small>Reconstruye texto como objetos PPTX cuando convenga.</small>
                <select id="editable-text-mode" name="editable_text_mode" aria-label="Modo de texto OCR">
                  <option value="large" selected>Solo títulos y texto grande</option>
                  <option value="all">Todo el texto detectado experimental</option>
                  <option value="none">Forzar texto raster exacto</option>
                </select>
              </div>
            </div>
            <div class="check">
              <input id="brain" name="brain" type="checkbox" value="true">
              <div>
                <label for="brain">Agentes OpenRouter</label>
                <small id="brain-note">Candidatos + juez free.</small>
                <div class="key-row" id="key-row" hidden>
                  <input id="openrouter-key" type="password" placeholder="sk-or-..." autocomplete="off">
                  <button class="small-button" id="save-key" type="button">Guardar</button>
                </div>
              </div>
            </div>
            <div class="check">
              <input id="compact" name="compact" type="checkbox" value="true">
              <div>
                <label for="compact">PPT más liviano</label>
                <small>Reduce objetos cuando el archivo pese demasiado.</small>
              </div>
            </div>
          </div>

          <div class="segmented" role="radiogroup" aria-label="Detalle">
            <input id="detail-heavy" type="radio" name="detail_mode" value="heavy" checked>
            <label for="detail-heavy">Pesado</label>
            <input id="detail-auto" type="radio" name="detail_mode" value="auto">
            <label for="detail-auto">Equilibrado</label>
            <input id="detail-grouped" type="radio" name="detail_mode" value="grouped">
            <label for="detail-grouped">Liviano</label>
          </div>

          <div class="segmented output-mode" role="radiogroup" aria-label="Salida">
            <input id="output-deck" type="radio" name="output_mode" value="deck" checked>
            <label for="output-deck">Deck único</label>
            <input id="output-zip" type="radio" name="output_mode" value="zip">
            <label for="output-zip">PPTX individuales ZIP</label>
          </div>

          <div class="actions">
            <button id="submit" type="submit">Generar PPTX</button>
            <a class="download" id="download" hidden>Descargar</a>
          </div>
        </form>
        </div>

        <div class="tab-panel models-panel" id="panel-models" hidden>
          <div>
            <label for="models-text">Modelos OpenRouter</label>
            <p class="drop-meta">Uno por línea o separados por coma. Se prueban en orden.</p>
          </div>
          <textarea id="models-text" spellcheck="false"></textarea>
          <div class="models-actions">
            <button id="save-models" type="button">Guardar modelos</button>
            <button id="reset-models" class="small-button" type="button">Free defaults</button>
            <span class="model-status" id="model-status">Cargando modelos.</span>
          </div>
        </div>
      </section>

      <section class="status" aria-label="Estado">
        <div class="meter-head">
          <span id="status-label">Esperando imagen</span>
          <span id="status-time">00:00</span>
        </div>
        <div class="meter" aria-label="Progreso">
          <div id="bar"></div>
        </div>
        <p class="stage" id="stage">Listo para procesar</p>
        <p class="message" id="message">El modo pesado divide flechas, iconos y detalles internos.</p>

        <div class="steps" id="steps">
          <div class="step" data-at="8"><span class="dot"></span><span>Subida</span></div>
          <div class="step" data-at="20"><span class="dot"></span><span>Análisis</span></div>
          <div class="step" data-at="58"><span class="dot"></span><span>Capas</span></div>
          <div class="step" data-at="76"><span class="dot"></span><span>PPTX</span></div>
          <div class="step" data-at="92"><span class="dot"></span><span>Calidad</span></div>
          <div class="step" data-at="100"><span class="dot"></span><span>Descarga</span></div>
        </div>

        <div class="metrics">
          <div class="metric"><span>Nota</span><strong id="grade">--</strong></div>
          <div class="metric"><span>MAE</span><strong id="mae">--</strong></div>
          <div class="metric"><span>Edge F1</span><strong id="edge">--</strong></div>
        </div>
        <p class="result-line" id="result-line">Sin job activo.</p>
      </section>
    </div>

    <footer>
      <div class="stat"><span>Tanda completa</span><strong>24/24</strong></div>
      <div class="stat"><span>MAE medio</span><strong>4.106</strong></div>
      <div class="stat"><span>Edge F1 medio</span><strong>0.864</strong></div>
      <div class="stat"><span>Modo recomendado</span><strong>Pesado</strong></div>
    </footer>
  </main>

  <script>
    window.OPENROUTER_AVAILABLE = {openrouter_available};

    const form = document.querySelector("#job-form");
    const fileInput = document.querySelector("#file");
    const filename = document.querySelector("#filename");
    const submit = document.querySelector("#submit");
    const download = document.querySelector("#download");
    const bar = document.querySelector("#bar");
    const statusLabel = document.querySelector("#status-label");
    const statusTime = document.querySelector("#status-time");
    const stage = document.querySelector("#stage");
    const message = document.querySelector("#message");
    const grade = document.querySelector("#grade");
    const mae = document.querySelector("#mae");
    const edge = document.querySelector("#edge");
    const resultLine = document.querySelector("#result-line");
    const brain = document.querySelector("#brain");
    const brainNote = document.querySelector("#brain-note");
    const keyRow = document.querySelector("#key-row");
    const openrouterKey = document.querySelector("#openrouter-key");
    const saveKey = document.querySelector("#save-key");
    const tabButtons = Array.from(document.querySelectorAll(".tab-button"));
    const panels = {{
      generate: document.querySelector("#panel-generate"),
      models: document.querySelector("#panel-models"),
    }};
    const modelsText = document.querySelector("#models-text");
    const saveModels = document.querySelector("#save-models");
    const resetModels = document.querySelector("#reset-models");
    const modelStatus = document.querySelector("#model-status");
    const steps = Array.from(document.querySelectorAll(".step"));

    let startedAt = null;
    let timer = null;

    if (!window.OPENROUTER_AVAILABLE) {{
      brain.disabled = true;
      keyRow.hidden = false;
      brainNote.textContent = "Guarda tu clave local para activarlo.";
    }}

    tabButtons.forEach((button) => {{
      button.addEventListener("click", () => activateTab(button.dataset.tab));
    }});

    loadModelSettings();

    saveKey.addEventListener("click", async () => {{
      const key = openrouterKey.value.trim();
      if (!key) return;
      saveKey.disabled = true;
      saveKey.textContent = "Guardando";
      try {{
        const data = new FormData();
        data.append("api_key", key);
        const response = await fetch("/settings/openrouter-key", {{ method: "POST", body: data }});
        if (!response.ok) throw new Error(await response.text());
        openrouterKey.value = "";
        keyRow.hidden = true;
        brain.disabled = false;
        brain.checked = true;
        brainNote.textContent = "Activo con modelos free.";
        window.OPENROUTER_AVAILABLE = true;
        if (Array.isArray((await response.clone().json()).models)) {{
          await loadModelSettings();
        }}
      }} catch (error) {{
        brainNote.textContent = error.message || String(error);
      }} finally {{
        saveKey.disabled = false;
        saveKey.textContent = "Guardar";
      }}
    }});

    saveModels.addEventListener("click", async () => {{
      await saveModelSettings(modelsText.value);
    }});

    resetModels.addEventListener("click", async () => {{
      const response = await fetch("/settings/openrouter");
      const settings = await response.json();
      modelsText.value = settings.default_models.join("\\n");
      await saveModelSettings(modelsText.value);
    }});

    function activateTab(name) {{
      tabButtons.forEach((button) => button.classList.toggle("active", button.dataset.tab === name));
      Object.entries(panels).forEach(([key, panel]) => {{
        panel.hidden = key !== name;
      }});
    }}

    async function loadModelSettings() {{
      try {{
        const response = await fetch("/settings/openrouter");
        if (!response.ok) throw new Error(await response.text());
        const settings = await response.json();
        modelsText.value = settings.models.join("\\n");
        modelStatus.textContent = settings.has_key ? "Clave activa." : "Sin clave guardada.";
      }} catch (error) {{
        modelStatus.textContent = error.message || String(error);
      }}
    }}

    async function saveModelSettings(value) {{
      const data = new FormData();
      data.append("models", value);
      saveModels.disabled = true;
      modelStatus.textContent = "Guardando modelos.";
      try {{
        const response = await fetch("/settings/openrouter-models", {{ method: "POST", body: data }});
        if (!response.ok) throw new Error(await response.text());
        const settings = await response.json();
        modelsText.value = settings.models.join("\\n");
        modelStatus.textContent = `${{settings.models.length}} modelos guardados.`;
      }} catch (error) {{
        modelStatus.textContent = error.message || String(error);
      }} finally {{
        saveModels.disabled = false;
      }}
    }}

    fileInput.addEventListener("change", () => {{
      const files = Array.from(fileInput.files);
      if (!files.length) {{
        filename.textContent = "Ningún archivo seleccionado";
        return;
      }}
      const totalSize = files.reduce((sum, file) => sum + file.size, 0);
      filename.textContent = files.length === 1
        ? `${{files[0].name}} · ${{formatBytes(totalSize)}}`
        : `${{files.length}} imágenes · ${{formatBytes(totalSize)}}`;
    }});

    form.addEventListener("submit", async (event) => {{
      event.preventDefault();
      if (!fileInput.files.length) return;
      resetUi();
      const data = new FormData(form);
      submit.disabled = true;
      submit.textContent = "Procesando";
      startedAt = Date.now();
      timer = window.setInterval(updateElapsed, 500);

      try {{
        const response = await fetch("/jobs", {{ method: "POST", body: data }});
        if (!response.ok) throw new Error(await response.text());
        const job = await response.json();
        poll(job.id);
      }} catch (error) {{
        showError(error.message || String(error));
      }}
    }});

    async function poll(id) {{
      try {{
        const response = await fetch(`/jobs/${{id}}`);
        if (!response.ok) throw new Error(await response.text());
        const job = await response.json();
        renderJob(job);
        if (job.status === "done" || job.status === "error") {{
          submit.disabled = false;
          submit.textContent = "Generar PPTX";
          window.clearInterval(timer);
          updateElapsed();
          return;
        }}
        window.setTimeout(() => poll(id), 700);
      }} catch (error) {{
        showError(error.message || String(error));
      }}
    }}

    function renderJob(job) {{
      const progress = Math.max(0, Math.min(100, Number(job.progress || 0)));
      bar.style.width = `${{progress}}%`;
      statusLabel.textContent = job.status === "done" ? "Completado" : job.status === "error" ? "Error" : `${{progress}}%`;
      stage.textContent = job.phase || "Procesando";
      message.textContent = job.message || "";
      resultLine.classList.toggle("error", job.status === "error");
      resultLine.textContent = job.error || job.summary || "Job activo.";
      steps.forEach((step) => step.classList.toggle("done", progress >= Number(step.dataset.at)));

      if (job.quality) {{
        grade.textContent = `${{job.quality.grade}} ${{Math.round(job.quality.score)}}`;
        mae.textContent = number(job.quality.metrics.mae, 2);
        edge.textContent = number(job.quality.metrics.edge_f1, 3);
      }}
      if (job.status === "done" && job.download_url) {{
        download.hidden = false;
        download.href = job.download_url;
        download.setAttribute("download", job.filename || "magical_layers_output.pptx");
      }}
      if (job.brain && job.brain.model_used) {{
        brainNote.textContent = `Cerebro: ${{job.brain.model_used}}`;
      }}
    }}

    function resetUi() {{
      download.hidden = true;
      download.removeAttribute("href");
      bar.style.width = "0%";
      statusLabel.textContent = "Iniciando";
      stage.textContent = "Preparando job";
      message.textContent = "Subiendo imágenes.";
      resultLine.classList.remove("error");
      resultLine.textContent = "Job activo.";
      grade.textContent = "--";
      mae.textContent = "--";
      edge.textContent = "--";
      steps.forEach((step) => step.classList.remove("done"));
    }}

    function showError(text) {{
      submit.disabled = false;
      submit.textContent = "Generar PPTX";
      window.clearInterval(timer);
      bar.style.width = "100%";
      statusLabel.textContent = "Error";
      stage.textContent = "No se pudo generar";
      message.textContent = "Revisa el archivo o vuelve a intentar con modo Auto.";
      resultLine.classList.add("error");
      resultLine.textContent = text;
    }}

    function updateElapsed() {{
      if (!startedAt) {{
        statusTime.textContent = "00:00";
        return;
      }}
      const seconds = Math.max(0, Math.round((Date.now() - startedAt) / 1000));
      const mm = String(Math.floor(seconds / 60)).padStart(2, "0");
      const ss = String(seconds % 60).padStart(2, "0");
      statusTime.textContent = `${{mm}}:${{ss}}`;
    }}

    function number(value, digits) {{
      const parsed = Number(value);
      return Number.isFinite(parsed) ? parsed.toFixed(digits) : "--";
    }}

    function formatBytes(bytes) {{
      if (!bytes) return "0 B";
      const units = ["B", "KB", "MB", "GB"];
      const index = Math.min(Math.floor(Math.log(bytes) / Math.log(1024)), units.length - 1);
      return `${{(bytes / Math.pow(1024, index)).toFixed(index ? 1 : 0)}} ${{units[index]}}`;
    }}
  </script>
</body>
</html>
"""


@app.post("/jobs")
async def create_job(
    files: list[UploadFile] | None = File(None),
    file: UploadFile | None = File(None),
    ocr: bool = Form(False),
    quality_check: bool = Form(False),
    brain: bool = Form(False),
    compact: bool = Form(False),
    detail_mode: str = Form("heavy"),
    editable_text_mode: str = Form("large"),
    output_mode: str = Form("deck"),
) -> dict[str, Any]:
    uploads = [upload for upload in (files or []) if upload.filename]
    if file is not None and file.filename:
        uploads.append(file)
    if not uploads:
        raise HTTPException(status_code=400, detail="Sube al menos una imagen.")

    job_id = uuid.uuid4().hex
    workdir = Path(tempfile.mkdtemp(prefix=f"magical_layers_{job_id}_"))
    input_paths: list[Path] = []
    for index, upload in enumerate(uploads, start=1):
        suffix = Path(upload.filename or "image.png").suffix.lower() or ".png"
        if suffix not in IMAGE_EXTENSIONS:
            raise HTTPException(status_code=400, detail=f"Formato no soportado: {upload.filename}")
        input_path = workdir / f"input_{index:03d}{suffix}"
        input_path.write_bytes(await upload.read())
        input_paths.append(input_path)

    total_files = len(input_paths)
    parsed_output_mode = _output_mode(output_mode, total_files)
    output_path = workdir / ("magical_layers_output.zip" if parsed_output_mode == "zip" else "magical_layers_output.pptx")

    JOBS[job_id] = {
        "id": job_id,
        "status": "queued",
        "progress": 3,
        "phase": "En cola",
        "message": f"{total_files} imagen{'es' if total_files != 1 else ''} recibida{'s' if total_files != 1 else ''}.",
        "summary": "Esperando turno local.",
        "created_at": time.time(),
        "filename": _download_filename([upload.filename for upload in uploads], parsed_output_mode),
        "media_type": ZIP_MEDIA_TYPE if parsed_output_mode == "zip" else PPTX_MEDIA_TYPE,
        "input_paths": [str(path) for path in input_paths],
        "output_path": str(output_path),
        "total_files": total_files,
        "quality": None,
        "download_url": None,
    }
    options = {
        "ocr": ocr,
        "quality_check": quality_check,
        "brain": brain and bool(os.getenv("OPENROUTER_API_KEY")),
        "compact": compact,
        "detail_mode": detail_mode,
        "editable_text_mode": _editable_text_mode(editable_text_mode),
        "output_mode": parsed_output_mode,
        "original_names": [upload.filename for upload in uploads],
    }
    asyncio.create_task(_run_job_async(job_id, input_paths, output_path, options))
    return _public_job(job_id)


@app.get("/jobs/{job_id}")
def get_job(job_id: str) -> dict[str, Any]:
    if job_id not in JOBS:
        raise HTTPException(status_code=404, detail="Job no encontrado.")
    return _public_job(job_id)


@app.get("/settings/openrouter")
def openrouter_settings() -> dict[str, Any]:
    return {
        "has_key": bool(os.getenv("OPENROUTER_API_KEY")),
        "models": judge_model_chain(),
        "default_models": DEFAULT_FREE_BRAIN_MODELS,
        "base_url": os.getenv("OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1"),
    }


@app.post("/settings/openrouter-key")
def save_openrouter_key(api_key: str = Form(...)) -> dict[str, Any]:
    cleaned = api_key.strip()
    if not cleaned.startswith("sk-or-"):
        raise HTTPException(status_code=400, detail="La clave OpenRouter debe empezar por sk-or-.")
    ENV_PATH.touch(exist_ok=True)
    set_key(str(ENV_PATH), "OPENROUTER_API_KEY", cleaned)
    if not os.getenv("OPENROUTER_BASE_URL"):
        set_key(str(ENV_PATH), "OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1")
    if not os.getenv("OPENROUTER_BRAIN_MODELS"):
        set_key(str(ENV_PATH), "OPENROUTER_BRAIN_MODELS", ",".join(DEFAULT_FREE_BRAIN_MODELS))
    if not os.getenv("OPENROUTER_JUDGE_MODELS"):
        set_key(str(ENV_PATH), "OPENROUTER_JUDGE_MODELS", ",".join(DEFAULT_FREE_BRAIN_MODELS))
    os.environ["OPENROUTER_API_KEY"] = cleaned
    os.environ.setdefault("OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1")
    os.environ.setdefault("OPENROUTER_BRAIN_MODELS", ",".join(DEFAULT_FREE_BRAIN_MODELS))
    os.environ.setdefault("OPENROUTER_JUDGE_MODELS", os.environ["OPENROUTER_BRAIN_MODELS"])
    return openrouter_settings() | {"ok": True}


@app.post("/settings/openrouter-models")
def save_openrouter_models(models: str = Form(...)) -> dict[str, Any]:
    parsed = _parse_models(models)
    if not parsed:
        raise HTTPException(status_code=400, detail="Añade al menos un modelo.")
    ENV_PATH.touch(exist_ok=True)
    value = ",".join(parsed)
    set_key(str(ENV_PATH), "OPENROUTER_BRAIN_MODELS", value)
    set_key(str(ENV_PATH), "OPENROUTER_JUDGE_MODELS", value)
    os.environ["OPENROUTER_BRAIN_MODELS"] = value
    os.environ["OPENROUTER_JUDGE_MODELS"] = value
    return openrouter_settings() | {"ok": True}


@app.get("/jobs/{job_id}/download")
def download_job(job_id: str) -> FileResponse:
    job = JOBS.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job no encontrado.")
    if job.get("status") != "done":
        raise HTTPException(status_code=409, detail="El PPTX aún no está listo.")
    output_path = Path(str(job["output_path"]))
    if not output_path.exists():
        raise HTTPException(status_code=404, detail="Archivo PPTX no encontrado.")
    return FileResponse(output_path, media_type=str(job.get("media_type") or PPTX_MEDIA_TYPE), filename=str(job["filename"]))


@app.post("/convert")
async def convert(file: UploadFile = File(...), ocr: bool = Form(False)) -> FileResponse:
    suffix = Path(file.filename or "image.png").suffix or ".png"
    workdir = Path(tempfile.mkdtemp(prefix="magical_layers_"))
    input_path = workdir / f"input{suffix}"
    output_path = workdir / "magical_layers_output.pptx"
    input_path.write_bytes(await file.read())

    await anyio.to_thread.run_sync(_convert_image, input_path, output_path, ocr)
    return FileResponse(output_path, media_type=PPTX_MEDIA_TYPE, filename="magical_layers_output.pptx")


async def _run_job_async(job_id: str, input_paths: list[Path], output_path: Path, options: dict[str, Any]) -> None:
    await anyio.to_thread.run_sync(_run_job, job_id, input_paths, output_path, options)


def _run_job(job_id: str, input_paths: list[Path], output_path: Path, options: dict[str, Any]) -> None:
    try:
        total = len(input_paths)
        _set_job(job_id, status="running", progress=8, phase="Subiendo", message="Archivos guardados en local.")
        time.sleep(0.08)
        _set_job(job_id, progress=20, phase="Analizando", message=_analysis_message(options))

        results = []
        total_layers = 0
        agent_rows: list[dict[str, Any]] = []
        judge_rows: list[dict[str, Any]] = []
        judge_errors: list[str] = []
        for index, input_path in enumerate(input_paths, start=1):
            start_progress = 20 + round((index - 1) / total * 42)
            end_progress = 20 + round(index / total * 42)
            _set_job(
                job_id,
                progress=start_progress,
                phase=f"Imagen {index}/{total}",
                message="Agentes OpenRouter generando candidatos." if options["brain"] else "Analizando y separando capas.",
            )
            if options["brain"]:
                selection = run_agentic_selection(
                    input_path,
                    output_path.parent,
                    index,
                    enable_ocr=bool(options["ocr"]),
                    editable_text_mode=str(options["editable_text_mode"]) if options["ocr"] else "none",
                    include_images=True,
                    ensemble=True,
                    progress=lambda message, current=index, count=total: _set_job(
                        job_id,
                        phase=f"Agentes {current}/{count}",
                        message=message,
                    ),
                )
                result = selection.result
                agent_rows.append(
                    {
                        "image_index": index,
                        "selected": selection.selected,
                        "candidates": _public_candidates(selection.candidates),
                        "judge_result": _compact_judge(selection.judge_result),
                        "judge_error": selection.judge_error,
                    }
                )
                if selection.judge_result:
                    judge_rows.append(selection.judge_result)
                if selection.judge_error:
                    judge_errors.append(selection.judge_error)
            else:
                pipeline_options = PipelineOptions(
                    enable_ocr=bool(options["ocr"]),
                    enable_brain=False,
                    editable_text_mode=str(options["editable_text_mode"]) if options["ocr"] else "none",
                    segment=_segment_options(str(options["detail_mode"]), bool(options["compact"])),
                )
                result = image_to_layers(input_path, pipeline_options)
            results.append(result)
            layer_count = int(result.metadata.get("raster_layers", len(result.layers)))
            total_layers += layer_count
            _set_job(
                job_id,
                progress=end_progress,
                phase=f"Capas {index}/{total}",
                message=f"{layer_count} capas raster detectadas.",
            )

        _set_job(
            job_id,
            progress=70,
            phase="Montando salida",
            message=f"{total} imagen{'es' if total != 1 else ''}, {total_layers} capas raster.",
        )
        pptx_paths: list[Path]
        if options.get("output_mode") == "zip":
            pptx_paths = _write_individual_pptx_zip(
                results,
                input_paths,
                output_path,
                [str(name or "") for name in options.get("original_names", [])],
            )
            shape_stats = _combined_shape_stats(pptx_paths)
        else:
            write_pptx_deck(results, output_path)
            pptx_paths = [output_path]
            shape_stats = _shape_stats(output_path)
        _set_job(
            job_id,
            progress=76,
            phase="Exportando PPTX",
            message=f"{shape_stats['shapes']} objetos en {shape_stats['slides']} slide{'s' if shape_stats['slides'] != 1 else ''}.",
        )

        quality = None
        if options["quality_check"]:
            _set_job(job_id, progress=84, phase="Midiendo calidad", message="Renderizando diapositivas.")
            if options.get("output_mode") == "zip":
                render_results = [
                    render_pptx_first_slide(path, output_path.parent / "renders" / f"{path.stem}.png") for path in pptx_paths
                ]
                metrics_rows = [
                    compare_render(input_path, render.image_path)
                    for input_path, render in zip(input_paths, render_results, strict=False)
                    if render.ok and render.image_path
                ]
            else:
                render_results = render_pptx_slides(output_path, output_path.parent / "renders", max_slides=total)
                metrics_rows = [
                    compare_render(input_path, render.image_path)
                    for input_path, render in zip(input_paths, render_results, strict=False)
                    if render.ok and render.image_path
                ]
            if metrics_rows:
                metrics = _average_metrics(metrics_rows)
                quality = _quality_from_metrics(metrics)
                quality["slides_measured"] = len(metrics_rows)
                _set_job(
                    job_id,
                    progress=94,
                    phase="Comparando render",
                    message=f"Nota {quality['grade']} media con MAE {metrics['mae']:.2f}.",
                    quality=quality,
                )
            else:
                error = render_results[0].error if render_results else "PowerPoint no pudo renderizar el deck."
                _set_job(
                    job_id,
                    progress=94,
                    phase="Calidad no disponible",
                    message=error or "PowerPoint no pudo renderizar el deck.",
                )
        else:
            _set_job(job_id, progress=94, phase="Saltando medición", message="PPTX generado sin comparación visual.")

        brain_state = _agent_state(agent_rows, judge_rows, judge_errors) or _brain_state([result.metadata for result in results])
        if brain_state:
            _set_job(job_id, brain=brain_state)
        summary = _summary_text([result.metadata for result in results], shape_stats, quality, bool(options["ocr"]), brain_state)
        _set_job(
            job_id,
            status="done",
            progress=100,
            phase="Salida lista",
            message="Descarga preparada.",
            summary=summary,
            download_url=f"/jobs/{job_id}/download",
            finished_at=time.time(),
        )
    except Exception as exc:
        _set_job(
            job_id,
            status="error",
            progress=100,
            phase="Error de proceso",
            message="La conversión se detuvo.",
            error=str(exc),
            finished_at=time.time(),
        )


def _convert_image(input_path: Path, output_path: Path, ocr: bool) -> None:
    options = PipelineOptions(enable_ocr=ocr, editable_text_mode="large" if ocr else "none", segment=segment_preset("heavy"))
    result = image_to_layers(input_path, options)
    write_pptx(result, output_path)


def _segment_options(detail_mode: str, compact: bool):
    return segment_preset(detail_mode, compact=compact)


def _editable_text_mode(value: str) -> str:
    return value if value in {"none", "large", "all"} else "large"


def _output_mode(value: str, total_files: int) -> str:
    if value == "zip" and total_files > 1:
        return "zip"
    return "deck"


def _write_individual_pptx_zip(
    results: list,
    input_paths: list[Path],
    output_path: Path,
    original_names: list[str],
) -> list[Path]:
    pptx_dir = output_path.parent / "individual_pptx"
    pptx_dir.mkdir(parents=True, exist_ok=True)
    pptx_paths: list[Path] = []
    manifest_rows = ["file,slides,objects,picture_layers,text_layers"]
    for index, (result, input_path) in enumerate(zip(results, input_paths, strict=False), start=1):
        original_name = original_names[index - 1] if index - 1 < len(original_names) else ""
        safe_stem = _safe_stem(Path(original_name) if original_name else input_path)
        pptx_path = pptx_dir / f"{index:02d}_{safe_stem}_layers.pptx"
        write_pptx(result, pptx_path)
        pptx_paths.append(pptx_path)
        stats = _shape_stats(pptx_path)
        manifest_rows.append(
            f"{pptx_path.name},{stats['slides']},{stats['shapes']},{stats['picture_shapes']},{stats['text_shapes']}"
        )
    with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for pptx_path in pptx_paths:
            archive.write(pptx_path, arcname=pptx_path.name)
        archive.writestr("LAYERS_MANIFEST.csv", "\n".join(manifest_rows) + "\n")
        archive.writestr("README_SHAREPOINT.txt", _sharepoint_readme())
    return pptx_paths


def _safe_stem(path: Path) -> str:
    stem = "".join(char if char.isalnum() else "_" for char in path.stem).strip("_")
    while "__" in stem:
        stem = stem.replace("__", "_")
    return stem or "image"


def _combined_shape_stats(pptx_paths: list[Path]) -> dict[str, int]:
    totals = {"slides": 0, "shapes": 0, "text_shapes": 0}
    for path in pptx_paths:
        stats = _shape_stats(path)
        totals["slides"] += stats["slides"]
        totals["shapes"] += stats["shapes"]
        totals["text_shapes"] += stats["text_shapes"]
    return totals


def _parse_models(value: str) -> list[str]:
    normalized = value.replace("\n", ",")
    models: list[str] = []
    for part in normalized.split(","):
        model = part.strip()
        if not model or model.startswith("#"):
            continue
        if model not in models:
            models.append(model)
    return models


def _shape_stats(pptx_path: Path) -> dict[str, int]:
    prs = Presentation(pptx_path)
    shape_count = 0
    picture_shapes = 0
    text_shapes = 0
    for slide in prs.slides:
        shape_count += len(slide.shapes)
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes += 1
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes += 1
    return {"slides": len(prs.slides), "shapes": shape_count, "picture_shapes": picture_shapes, "text_shapes": text_shapes}


def _sharepoint_readme() -> str:
    return (
        "Magical Layers - lectura en SharePoint / PowerPoint\n\n"
        "1. Extrae este ZIP antes de subirlo. SharePoint no edita las capas dentro de un ZIP.\n"
        "2. Sube o abre cada archivo .pptx individual, no el ZIP completo.\n"
        "3. En PowerPoint para la web, entra en modo edicion y abre Organizar > Panel de seleccion.\n"
        "4. Las capas raster estan nombradas como raster_001, raster_002, etc.\n"
        "5. Si el navegador lo muestra como una imagen plana o va lento, usa Abrir en la aplicacion de escritorio.\n"
        "6. LAYERS_MANIFEST.csv contiene el recuento de objetos y capas de imagen por PPTX.\n"
    )


def _average_metrics(rows: list[dict[str, float]]) -> dict[str, float]:
    keys = sorted({key for row in rows for key in row})
    return {key: sum(float(row.get(key, 0.0)) for row in rows) / len(rows) for key in keys}


def _quality_from_metrics(metrics: dict[str, float]) -> dict[str, Any]:
    mae = float(metrics.get("mae", 999.0))
    edge_f1 = float(metrics.get("edge_f1", 0.0))
    score = max(0.0, min(100.0, 100.0 - mae * 1.8 + (edge_f1 - 0.80) * 25.0))
    if score >= 92:
        grade = "A"
    elif score >= 86:
        grade = "B+"
    elif score >= 78:
        grade = "B"
    elif score >= 68:
        grade = "C"
    else:
        grade = "D"
    return {"grade": grade, "score": score, "metrics": metrics}


def _summary_text(
    metadata: list[dict[str, Any]],
    shape_stats: dict[str, int],
    quality: dict[str, Any] | None,
    ocr_enabled: bool,
    brain: dict[str, Any] | None = None,
) -> str:
    raster_layers = sum(int(row.get("raster_layers", 0)) for row in metadata)
    editable = sum(int(row.get("editable_text_lines", 0)) for row in metadata) if ocr_enabled else 0
    base = f"{shape_stats['slides']} slide{'s' if shape_stats['slides'] != 1 else ''}, {raster_layers} capas raster, {shape_stats['shapes']} objetos PPTX"
    if editable:
        base += f", {editable} textos editables"
    if quality:
        base += f", calidad {quality['grade']} ({quality['score']:.0f}/100)"
    if brain and brain.get("model_used"):
        base += f", cerebro {brain['model_used']}"
    return base + "."


def _brain_state(metadata: list[dict[str, Any]]) -> dict[str, Any] | None:
    plans = [row.get("brain_plan") for row in metadata if row.get("brain_plan")]
    if not plans:
        return None
    first = plans[0]
    if not isinstance(first, dict):
        return {"content": str(first)}
    state = {
        "model_used": first.get("model_used"),
        "json": first.get("json"),
        "content": first.get("content"),
    }
    if first.get("error"):
        state["error"] = first.get("error")
    if len(plans) > 1:
        state["slides_with_brain"] = len(plans)
    return {key: value for key, value in state.items() if value}


def _agent_state(
    agent_rows: list[dict[str, Any]],
    judge_rows: list[dict[str, Any]],
    judge_errors: list[str],
) -> dict[str, Any] | None:
    if not agent_rows:
        return None
    selected_modes = [str(row.get("selected", {}).get("candidate_mode", "")) for row in agent_rows]
    model_used = None
    for row in judge_rows:
        if row.get("model_used"):
            model_used = row.get("model_used")
            break
    return {
        "mode": "agentic_candidate_race",
        "model_used": model_used,
        "slides_judged": len(judge_rows),
        "slides_local_fallback": len(agent_rows) - len(judge_rows),
        "selected_modes": selected_modes,
        "judge_errors": judge_errors[:3],
        "runs": agent_rows,
    }


def _public_candidates(candidates: list[dict[str, Any]]) -> list[dict[str, Any]]:
    keys = [
        "candidate_mode",
        "candidate_preset",
        "candidate_compact",
        "raster_layers",
        "shapes",
        "mae",
        "edge_f1",
        "largest_layer_ratio",
        "internal_detail_layers",
        "local_score",
        "render_error",
    ]
    return [{key: candidate.get(key) for key in keys if key in candidate} for candidate in candidates]


def _compact_judge(judge_result: dict[str, Any] | None) -> dict[str, Any] | None:
    if not judge_result:
        return None
    keys = [
        "winner_index",
        "winner_reason",
        "quality_score",
        "failure_modes",
        "next_pipeline_change",
        "model_used",
        "input_mode",
        "model_errors",
    ]
    return {key: judge_result.get(key) for key in keys if key in judge_result}


def _analysis_message(options: dict[str, Any]) -> str:
    parts = ["segmentación pesada" if options.get("detail_mode") == "heavy" else "segmentación"]
    if options["ocr"]:
        text_mode = "OCR completo" if options.get("editable_text_mode") == "all" else "OCR"
        parts.append(text_mode)
    if options["brain"]:
        parts.append("agentes OpenRouter")
    return " + ".join(parts).capitalize() + " en marcha."


def _download_filename(original: str | list[str | None] | None, output_mode: str) -> str:
    if output_mode == "zip":
        if isinstance(original, list) and len(original) > 1:
            return f"magical_layers_{len(original)}_pptx_layers.zip"
        return "magical_layers_pptx_layers.zip"
    return _pptx_filename(original)


def _pptx_filename(original: str | list[str | None] | None) -> str:
    if isinstance(original, list):
        if len(original) > 1:
            return f"magical_layers_{len(original)}_slides_layers.pptx"
        original = original[0] if original else None
    stem = Path(original or "magical_layers_output").stem
    keep = []
    for char in stem:
        if char.isalnum() or char in "-_":
            keep.append(char)
        elif char.isspace():
            keep.append("_")
    safe = "".join(keep)[:80] or "magical_layers_output"
    return f"{safe}_layers.pptx"


def _set_job(job_id: str, **values: Any) -> None:
    job = JOBS[job_id]
    job.update(values)
    job["updated_at"] = time.time()


def _public_job(job_id: str) -> dict[str, Any]:
    job = JOBS[job_id]
    keys = [
        "id",
        "status",
        "progress",
        "phase",
        "message",
        "summary",
        "error",
        "quality",
        "brain",
        "download_url",
        "filename",
        "total_files",
        "created_at",
        "updated_at",
        "finished_at",
    ]
    return {key: job.get(key) for key in keys if key in job}
