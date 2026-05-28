from __future__ import annotations

from io import BytesIO
from pathlib import Path
from collections.abc import Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from .models import Layer, LayerKind, PipelineResult


EMU_PER_INCH = 914400


def write_pptx(result: PipelineResult, output_path: str | Path, dpi: float = 112.0) -> Path:
    return write_pptx_deck([result], output_path, dpi=dpi)


def write_pptx_deck(results: Sequence[PipelineResult], output_path: str | Path, dpi: float = 112.0) -> Path:
    if not results:
        raise ValueError("At least one pipeline result is required")

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    first = results[0]
    prs.slide_width = Emu(round(first.image_width / dpi * EMU_PER_INCH))
    prs.slide_height = Emu(round(first.image_height / dpi * EMU_PER_INCH))
    blank_layout = prs.slide_layouts[6]

    for result in results:
        _add_result_slide(prs, blank_layout, result)

    prs.save(output)
    return output


def _add_result_slide(prs: Presentation, blank_layout, result: PipelineResult) -> None:
    slide = prs.slides.add_slide(blank_layout)
    _paint_background(slide, result.background_color)

    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    scale = min(slide_width / result.image_width, slide_height / result.image_height)
    offset_x = (slide_width - result.image_width * scale) / 2
    offset_y = (slide_height - result.image_height * scale) / 2

    raster_layers = [layer for layer in result.layers if layer.kind == LayerKind.RASTER]
    text_layers = [layer for layer in result.layers if layer.kind == LayerKind.TEXT]

    for layer in raster_layers:
        if layer.image is None:
            continue
        stream = BytesIO()
        layer.image.save(stream, format="PNG")
        stream.seek(0)
        picture = slide.shapes.add_picture(
            stream,
            Emu(round(offset_x + layer.bbox.x * scale)),
            Emu(round(offset_y + layer.bbox.y * scale)),
            width=Emu(round(layer.bbox.width * scale)),
            height=Emu(round(layer.bbox.height * scale)),
        )
        picture.name = layer.id

    for layer in text_layers:
        _add_text_layer(slide, layer, scale, offset_x, offset_y)


def _paint_background(slide, color: tuple[int, int, int]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def _add_text_layer(slide, layer: Layer, scale: float, offset_x: float, offset_y: float) -> None:
    shape = slide.shapes.add_textbox(
        Emu(round(offset_x + layer.bbox.x * scale)),
        Emu(round(offset_y + layer.bbox.y * scale)),
        Emu(round(layer.bbox.width * scale)),
        Emu(round(layer.bbox.height * scale * 1.15)),
    )
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER if layer.align == "center" else PP_ALIGN.LEFT
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)

    runs = layer.runs if layer.runs else []
    if not runs and layer.text:
        from .models import TextRun

        runs = [TextRun(layer.text, layer.color or (0, 0, 0))]

    for text_run in runs:
        run = paragraph.add_run()
        run.text = text_run.text
        font = run.font
        font.name = "Aptos Display"
        font.size = Pt(layer.font_size or 18)
        font.bold = layer.bold
        font.color.rgb = RGBColor(*text_run.color)
