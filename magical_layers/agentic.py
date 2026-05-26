from __future__ import annotations

import json
import math
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable

from pptx import Presentation

from .evaluation import compare_render, make_comparison_image, render_pptx_first_slide
from .models import PipelineResult
from .openrouter_judge import JudgeConfig, judge_candidates, judge_model_chain
from .orchestrator import PipelineOptions, image_to_layers
from .pptx_writer import write_pptx
from .presets import segment_preset


ProgressCallback = Callable[[str], None]


@dataclass(slots=True)
class CandidateSpec:
    name: str
    preset: str
    compact: bool = False


@dataclass(slots=True)
class AgenticSelection:
    result: PipelineResult
    selected: dict[str, Any]
    candidates: list[dict[str, Any]]
    judge_result: dict[str, Any] | None = None
    judge_error: str | None = None


DEFAULT_CANDIDATES = [
    CandidateSpec("heavy", "heavy", compact=False),
    CandidateSpec("heavy_lite", "heavy", compact=True),
    CandidateSpec("granular", "granular", compact=False),
    CandidateSpec("balanced", "auto", compact=False),
    CandidateSpec("light", "grouped", compact=False),
]


def run_agentic_selection(
    image_path: Path,
    workdir: Path,
    image_index: int,
    *,
    enable_ocr: bool,
    ocr_lang: str = "es",
    editable_text_mode: str = "large",
    editable_text_min_height: int | None = None,
    include_images: bool = True,
    ensemble: bool = True,
    progress: ProgressCallback | None = None,
) -> AgenticSelection:
    candidate_dir = workdir / "agent_candidates" / f"image_{image_index:03d}"
    candidate_dir.mkdir(parents=True, exist_ok=True)
    candidates: list[dict[str, Any]] = []
    candidate_results: list[PipelineResult] = []

    for spec_index, spec in enumerate(DEFAULT_CANDIDATES):
        if progress:
            progress(f"candidato {spec_index + 1}/{len(DEFAULT_CANDIDATES)}: {spec.name}")
        candidate, result = _process_candidate(
            image_path=image_path,
            candidate_dir=candidate_dir,
            spec=spec,
            image_index=image_index,
            enable_ocr=enable_ocr,
            ocr_lang=ocr_lang,
            editable_text_mode=editable_text_mode,
            editable_text_min_height=editable_text_min_height,
        )
        candidates.append(candidate)
        candidate_results.append(result)

    judge_result: dict[str, Any] | None = None
    judge_error: str | None = None
    if os.getenv("OPENROUTER_API_KEY"):
        if progress:
            progress("juzgando candidatos con OpenRouter")
        judge_result, judge_error = _judge_candidates(image_path.name, candidates, include_images, ensemble)

    selected_index = _selected_index(candidates, judge_result)
    selected = dict(candidates[selected_index])
    selected["selected_by"] = "openrouter" if judge_result else "local_score"
    selected["selected_index"] = selected_index
    return AgenticSelection(
        result=candidate_results[selected_index],
        selected=selected,
        candidates=candidates,
        judge_result=judge_result,
        judge_error=judge_error,
    )


def _process_candidate(
    image_path: Path,
    candidate_dir: Path,
    spec: CandidateSpec,
    image_index: int,
    enable_ocr: bool,
    ocr_lang: str,
    editable_text_mode: str,
    editable_text_min_height: int | None,
) -> tuple[dict[str, Any], PipelineResult]:
    stem = _safe_stem(image_path.stem)
    prefix = f"{image_index:03d}_{stem}_{spec.name}"
    pptx_path = candidate_dir / f"{prefix}.pptx"
    render_path = candidate_dir / f"{prefix}.png"
    preview_path = candidate_dir / f"{prefix}.jpg"
    options = PipelineOptions(
        enable_ocr=enable_ocr,
        ocr_lang=ocr_lang,
        editable_text_mode=editable_text_mode if enable_ocr else "none",
        editable_text_min_height=editable_text_min_height,
        segment=segment_preset(spec.preset, compact=spec.compact),
    )
    result = image_to_layers(image_path, options)
    write_pptx(result, pptx_path)
    shape_stats = _shape_stats(pptx_path)
    row: dict[str, Any] = {
        "candidate_mode": spec.name,
        "candidate_preset": spec.preset,
        "candidate_compact": spec.compact,
        "pptx": str(pptx_path),
        "ocr_lines": result.metadata.get("ocr_lines", 0),
        "editable_text_lines": result.metadata.get("editable_text_lines", 0),
        "raster_layers": result.metadata.get("raster_layers", 0),
        "shapes": shape_stats["shapes"],
        "text_shapes": shape_stats["text_shapes"],
    }
    row.update(result.metadata)
    row.update(_layer_quality_features(result))
    render = render_pptx_first_slide(pptx_path, render_path)
    if render.ok and render.image_path:
        metrics = compare_render(image_path, render.image_path)
        row.update(metrics)
        make_comparison_image(
            image_path,
            render.image_path,
            preview_path,
            title=(
                f"{image_path.name} | {spec.name} | layers {row['raster_layers']} | "
                f"MAE {metrics['mae']:.2f} | Edge F1 {metrics['edge_f1']:.3f}"
            ),
        )
        row["render"] = str(render.image_path)
        row["preview"] = str(preview_path)
    else:
        row["render_error"] = render.error or "unknown render error"
    row["local_score"] = _local_score(row)
    return row, result


def _judge_candidates(
    image_name: str,
    candidates: list[dict[str, Any]],
    include_images: bool,
    ensemble: bool,
) -> tuple[dict[str, Any] | None, str | None]:
    if not ensemble:
        models = judge_model_chain()
        config = JudgeConfig(
            enabled=True,
            model=models[0],
            fallback_models=models[1:],
            include_images=include_images,
        )
        try:
            return judge_candidates(image_name, candidates, config), None
        except Exception as exc:
            return None, str(exc)

    results: list[dict[str, Any]] = []
    errors: dict[str, str] = {}
    models = judge_model_chain()
    for model in models:
        config = JudgeConfig(enabled=True, model=model, fallback_models=[], include_images=include_images)
        try:
            result = judge_candidates(image_name, candidates, config)
            if result:
                results.append(result)
        except Exception as exc:
            errors[model] = str(exc)

    if not results:
        return None, json.dumps(errors, ensure_ascii=False) if errors else "no judge result"

    votes: dict[int, int] = {}
    for result in results:
        try:
            index = int(result.get("winner_index", -1))
        except Exception:
            continue
        if 0 <= index < len(candidates):
            votes[index] = votes.get(index, 0) + 1

    if votes:
        local_best = _selected_index(candidates, None)
        winner_index = max(votes, key=lambda index: (votes[index], index == local_best))
    else:
        winner_index = _selected_index(candidates, None)

    return (
        {
            "winner_index": winner_index,
            "winner_reason": "Ensemble vote across free OpenRouter judge models.",
            "quality_score": _average_quality(results),
            "failure_modes": _collect_failure_modes(results),
            "next_pipeline_change": _collect_next_change(results),
            "model_used": "ensemble:" + ",".join(str(result.get("model_used", "")) for result in results),
            "input_mode": ",".join(sorted({str(result.get("input_mode", "")) for result in results if result.get("input_mode")})),
            "model_results": results,
            "model_errors": errors,
        },
        json.dumps(errors, ensure_ascii=False) if errors else None,
    )


def _selected_index(candidates: list[dict[str, Any]], judge_result: dict[str, Any] | None) -> int:
    if judge_result and "winner_index" in judge_result:
        try:
            index = int(judge_result["winner_index"])
            if 0 <= index < len(candidates):
                return index
        except Exception:
            pass
    return max(range(len(candidates)), key=lambda index: float(candidates[index].get("local_score", -9999)))


def _local_score(row: dict[str, Any]) -> float:
    mae = float(row.get("mae", 999.0))
    edge_f1 = float(row.get("edge_f1", 0.0))
    layers = max(1.0, float(row.get("raster_layers", 1)))
    largest_ratio = float(row.get("largest_layer_ratio", 1.0))
    internal_details = float(row.get("internal_detail_layers", 0))
    visual = -mae * 0.35 + edge_f1 * 4.0
    editability = min(2.6, math.log1p(layers) / math.log1p(1200) * 2.6)
    detail_bonus = min(0.8, internal_details / 250.0)
    grouping_penalty = max(0.0, largest_ratio - 0.22) * 3.0
    return visual + editability + detail_bonus - grouping_penalty


def _layer_quality_features(result: PipelineResult) -> dict[str, float | int]:
    raster_layers = [layer for layer in result.layers if layer.image is not None]
    if not raster_layers:
        return {"largest_layer_ratio": 1.0, "internal_detail_layers": 0, "internal_detail_base_layers": 0}
    slide_area = max(1, result.image_width * result.image_height)
    largest_area = max(layer.bbox.area for layer in raster_layers)
    internal_details = sum(1 for layer in raster_layers if layer.metadata.get("internal_detail"))
    internal_bases = sum(1 for layer in raster_layers if layer.metadata.get("internal_detail_base"))
    return {
        "largest_layer_ratio": largest_area / slide_area,
        "internal_detail_layers": internal_details,
        "internal_detail_base_layers": internal_bases,
    }


def _shape_stats(pptx_path: Path) -> dict[str, int]:
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    text_shapes = 0
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            text_shapes += 1
    return {"shapes": len(slide.shapes), "text_shapes": text_shapes}


def _average_quality(results: list[dict[str, Any]]) -> float:
    scores: list[float] = []
    for result in results:
        try:
            score = float(result.get("quality_score", 0))
            if score > 1:
                score = score / 100
            scores.append(max(0.0, min(1.0, score)))
        except Exception:
            continue
    return sum(scores) / len(scores) if scores else 0.0


def _collect_failure_modes(results: list[dict[str, Any]]) -> list[str]:
    modes: list[str] = []
    for result in results:
        value = result.get("failure_modes", [])
        if isinstance(value, list):
            for item in value:
                text = str(item)
                if text and text not in modes:
                    modes.append(text)
    return modes[:8]


def _collect_next_change(results: list[dict[str, Any]]) -> str:
    changes: list[str] = []
    for result in results:
        text = str(result.get("next_pipeline_change", "")).strip()
        if text and text not in changes:
            changes.append(text)
    return " | ".join(changes[:4])


def _safe_stem(stem: str) -> str:
    keep = []
    for char in stem:
        if char.isalnum() or char in "-_":
            keep.append(char)
        elif char.isspace():
            keep.append("_")
    return "".join(keep)[:80] or "image"
