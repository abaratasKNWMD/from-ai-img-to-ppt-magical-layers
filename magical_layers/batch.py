from __future__ import annotations

import argparse
import csv
import json
import os
import shutil
from pathlib import Path

from dotenv import load_dotenv
from pptx import Presentation

from .evaluation import compare_render, make_comparison_image, render_pptx_first_slide, write_html_report
from .openrouter_judge import DEFAULT_JUDGE_MODELS, JudgeConfig, judge_candidates
from .orchestrator import PipelineOptions, image_to_layers
from .pptx_writer import write_pptx
from .presets import segment_preset
from .segmentation import SegmentOptions


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Batch-convert images to layered PPTX and score rendered output.")
    parser.add_argument("--input-dir", type=Path, default=Path("images"))
    parser.add_argument("--output-dir", type=Path, default=Path("outputs/batch"))
    parser.add_argument("--ocr-lang", default="es")
    parser.add_argument("--limit", type=int, default=None)
    parser.add_argument("--start", type=int, default=1, help="1-based start index in sorted image list.")
    parser.add_argument("--no-render", action="store_true")
    parser.add_argument(
        "--preset",
        choices=["heavy", "auto", "granular", "grouped", "custom"],
        default="heavy",
        help="Segmentation preset. Use custom to honor low-level threshold/min-area/dilation flags.",
    )
    parser.add_argument("--bg-threshold", type=float, default=35.0)
    parser.add_argument("--bg-thresholds", default=None, help="Comma-separated thresholds for auto mode, for example 35,18.")
    parser.add_argument("--surface-threshold", type=float, default=18.0)
    parser.add_argument("--dilation", type=int, default=3)
    parser.add_argument("--min-area", type=int, default=100)
    parser.add_argument("--max-layers", type=int, default=500)
    parser.add_argument("--dpi", type=float, default=112.0)
    parser.add_argument("--editable-text", choices=["auto", "none", "large", "all"], default="auto")
    parser.add_argument("--editable-min-height", type=int, default=None)
    parser.add_argument("--judge", action="store_true", help="Ask OpenRouter to judge candidate quality.")
    parser.add_argument("--judge-model", default=None)
    parser.add_argument("--judge-models", default=None, help="Comma-separated OpenRouter judge model fallback list.")
    parser.add_argument("--judge-ensemble", action="store_true", help="Call every judge model and aggregate votes.")
    parser.add_argument("--judge-text-only", action="store_true", help="Do not send preview images to OpenRouter.")
    return parser


def main() -> None:
    load_dotenv()
    args = build_parser().parse_args()
    if args.judge and not os.getenv("OPENROUTER_API_KEY"):
        raise SystemExit("OPENROUTER_API_KEY is not set. Put it in .env or export it before using --judge.")
    paths = [p for p in sorted(args.input_dir.iterdir()) if p.suffix.lower() in IMAGE_EXTENSIONS]
    start = max(0, args.start - 1)
    paths = paths[start:]
    if args.limit is not None:
        paths = paths[: args.limit]

    args.output_dir.mkdir(parents=True, exist_ok=True)
    preview_dir = args.output_dir / "previews"
    render_dir = args.output_dir / "renders"
    pptx_dir = args.output_dir / "pptx"
    meta_dir = args.output_dir / "metadata"
    for directory in (preview_dir, render_dir, pptx_dir, meta_dir):
        directory.mkdir(parents=True, exist_ok=True)

    rows: list[dict[str, object]] = []
    for offset, image_path in enumerate(paths, start=args.start):
        safe_stem = _safe_stem(image_path.stem)
        print(f"[{offset}/{args.start + len(paths) - 1}] {image_path.name}", flush=True)
        metadata_path = meta_dir / f"{offset:02d}_{safe_stem}.json"
        modes = ["none", "large"] if args.editable_text == "auto" else [args.editable_text]
        thresholds = _thresholds(args)
        candidates = [
            _process_candidate(args, image_path, offset, safe_stem, mode, threshold, pptx_dir, render_dir, preview_dir)
            for mode in modes
            for threshold in thresholds
        ]
        judge_result, judge_error = _judge_if_enabled(args, image_path.name, candidates)
        row = _choose_candidate(candidates, judge_result)
        row["index"] = offset
        row["name"] = image_path.name
        row["selected_mode"] = row.get("candidate_mode", "")
        row["selected_threshold"] = row.get("candidate_threshold", "")
        if judge_result:
            row["judge_model"] = judge_result.get("model_used", args.judge_model or JudgeConfig().model)
            row["judge_result"] = json.dumps(judge_result, ensure_ascii=False)
        if judge_error:
            row["judge_error"] = judge_error

        final_pptx = pptx_dir / f"{offset:02d}_{safe_stem}.pptx"
        final_render = render_dir / f"{offset:02d}_{safe_stem}.png"
        final_preview = preview_dir / f"{offset:02d}_{safe_stem}.jpg"
        _copy_if_present(row.get("pptx"), final_pptx)
        _copy_if_present(row.get("render"), final_render)
        _copy_if_present(row.get("preview"), final_preview)
        row["pptx"] = str(final_pptx)
        if final_render.exists():
            row["render"] = str(final_render)
        if final_preview.exists():
            row["preview"] = str(final_preview)

        for candidate in candidates:
            mode = candidate.get("candidate_mode")
            threshold = candidate.get("candidate_threshold")
            if mode and threshold:
                key = f"{mode}_t{_threshold_label(float(threshold))}"
                row[f"mae_{key}"] = candidate.get("mae", "")
                row[f"edge_f1_{key}"] = candidate.get("edge_f1", "")

        metadata_path.write_text(json.dumps(row, indent=2, ensure_ascii=False), encoding="utf-8")
        rows.append(row)

    _write_csv(rows, args.output_dir / "report.csv")
    write_html_report(rows, args.output_dir / "report.html")
    print(f"Wrote {args.output_dir / 'report.html'}", flush=True)


def _process_candidate(
    args: argparse.Namespace,
    image_path: Path,
    index: int,
    safe_stem: str,
    mode: str,
    threshold: float,
    pptx_dir: Path,
    render_dir: Path,
    preview_dir: Path,
) -> dict[str, object]:
    suffix = "" if args.editable_text != "auto" and len(_thresholds(args)) == 1 else f"_{mode}_t{_threshold_label(threshold)}"
    pptx_path = pptx_dir / f"{index:02d}_{safe_stem}{suffix}.pptx"
    rendered_path = render_dir / f"{index:02d}_{safe_stem}{suffix}.png"
    preview_path = preview_dir / f"{index:02d}_{safe_stem}{suffix}.jpg"
    options = PipelineOptions(
        enable_ocr=mode != "none",
        ocr_lang=args.ocr_lang,
        editable_text_mode=mode,
        editable_text_min_height=args.editable_min_height,
        segment=_segment_options(args, threshold),
    )
    result = image_to_layers(image_path, options)
    write_pptx(result, pptx_path, dpi=args.dpi)
    shape_stats = _shape_stats(pptx_path)
    row: dict[str, object] = {
        "candidate_mode": mode,
        "candidate_threshold": threshold,
        "pptx": str(pptx_path),
        "ocr_lines": result.metadata.get("ocr_lines", 0),
        "editable_text_lines": result.metadata.get("editable_text_lines", 0),
        "raster_layers": result.metadata.get("raster_layers", 0),
        "shapes": shape_stats["shapes"],
        "text_shapes": shape_stats["text_shapes"],
    }
    row.update(result.metadata)
    if not args.no_render:
        render = render_pptx_first_slide(pptx_path, rendered_path)
        if render.ok and render.image_path:
            metrics = compare_render(image_path, render.image_path)
            row.update(metrics)
            make_comparison_image(
                image_path,
                render.image_path,
                preview_path,
                title=f"{image_path.name} | {mode} | MAE {metrics['mae']:.2f} | Edge F1 {metrics['edge_f1']:.3f}",
            )
            row["render"] = str(render.image_path)
            row["preview"] = str(preview_path)
        else:
            row["render_error"] = render.error or "unknown render error"
    return row


def _segment_options(args: argparse.Namespace, threshold: float) -> SegmentOptions:
    if args.preset != "custom":
        options = segment_preset(args.preset)
        options.bg_threshold = threshold
        return options
    return SegmentOptions(
        bg_threshold=threshold,
        surface_threshold=args.surface_threshold if threshold > args.surface_threshold else None,
        dilation_px=args.dilation,
        min_area_px=args.min_area,
        max_layers=args.max_layers,
    )


def _judge_if_enabled(
    args: argparse.Namespace,
    image_name: str,
    candidates: list[dict[str, object]],
) -> tuple[dict[str, object] | None, str | None]:
    if not args.judge:
        return None, None
    models = _judge_models(args)
    if args.judge_ensemble:
        return _judge_ensemble(args, image_name, candidates, models)
    config = JudgeConfig(
        enabled=True,
        model=models[0],
        fallback_models=models[1:],
        include_images=not args.judge_text_only,
    )
    try:
        return judge_candidates(image_name, candidates, config), None
    except Exception as exc:
        return None, str(exc)


def _judge_ensemble(
    args: argparse.Namespace,
    image_name: str,
    candidates: list[dict[str, object]],
    models: list[str],
) -> tuple[dict[str, object] | None, str | None]:
    results: list[dict[str, object]] = []
    errors: dict[str, str] = {}
    for model in models:
        config = JudgeConfig(
            enabled=True,
            model=model,
            fallback_models=[],
            include_images=not args.judge_text_only,
        )
        try:
            result = judge_candidates(image_name, candidates, config)
            if result is not None:
                results.append(result)
        except Exception as exc:
            errors[model] = str(exc)

    if not results:
        return None, json.dumps(errors, ensure_ascii=False)

    votes: dict[int, int] = {}
    for result in results:
        try:
            winner = int(result.get("winner_index", -1))
        except Exception:
            continue
        if 0 <= winner < len(candidates):
            votes[winner] = votes.get(winner, 0) + 1

    if votes:
        metric_best = _choose_candidate(candidates)
        metric_index = candidates.index(metric_best)
        winner_index = max(votes, key=lambda index: (votes[index], index == metric_index))
    else:
        winner_index = candidates.index(_choose_candidate(candidates))

    ensemble = {
        "winner_index": winner_index,
        "winner_reason": "Ensemble vote across OpenRouter judge models.",
        "quality_score": _average_quality(results),
        "failure_modes": _collect_failure_modes(results),
        "next_pipeline_change": _collect_next_change(results),
        "model_used": "ensemble",
        "input_mode": ",".join(sorted({str(result.get("input_mode", "")) for result in results if result.get("input_mode")})),
        "model_results": results,
        "model_errors": errors,
    }
    return ensemble, None if not errors else json.dumps(errors, ensure_ascii=False)


def _choose_candidate(
    candidates: list[dict[str, object]],
    judge_result: dict[str, object] | None = None,
) -> dict[str, object]:
    if judge_result and "winner_index" in judge_result:
        try:
            index = int(judge_result["winner_index"])
            if 0 <= index < len(candidates) and "mae" in candidates[index]:
                return candidates[index]
        except Exception:
            pass
    scored = [candidate for candidate in candidates if "mae" in candidate]
    if not scored:
        return candidates[0]
    return min(scored, key=lambda row: (float(row.get("mae", 9999)), -float(row.get("edge_f1", 0))))


def _copy_if_present(source: object, target: Path) -> None:
    if not source:
        return
    src = Path(str(source))
    if src.exists() and src.resolve() != target.resolve():
        shutil.copy2(src, target)


def _shape_stats(pptx_path: Path) -> dict[str, int]:
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    text_shapes = 0
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            text_shapes += 1
    return {"shapes": len(slide.shapes), "text_shapes": text_shapes}


def _write_csv(rows: list[dict[str, object]], output: Path) -> None:
    keys: list[str] = []
    for row in rows:
        for key in row:
            if key not in keys:
                keys.append(key)
    with output.open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=keys)
        writer.writeheader()
        writer.writerows(rows)


def _safe_stem(stem: str) -> str:
    keep = []
    for char in stem:
        if char.isalnum() or char in "-_":
            keep.append(char)
        elif char.isspace():
            keep.append("_")
    return "".join(keep)[:80] or "image"


def _thresholds(args: argparse.Namespace) -> list[float]:
    if args.bg_thresholds:
        values = [float(part.strip()) for part in args.bg_thresholds.split(",") if part.strip()]
    elif args.editable_text == "auto":
        values = [args.bg_threshold, 18.0]
    else:
        values = [args.bg_threshold]
    deduped: list[float] = []
    for value in values:
        if value not in deduped:
            deduped.append(value)
    return deduped


def _threshold_label(value: float) -> str:
    if value.is_integer():
        return str(int(value))
    return str(value).replace(".", "p")


def _judge_models(args: argparse.Namespace) -> list[str]:
    if args.judge_models:
        return [part.strip() for part in args.judge_models.split(",") if part.strip()]
    if args.judge_model:
        return [args.judge_model, *[model for model in DEFAULT_JUDGE_MODELS if model != args.judge_model]]
    return DEFAULT_JUDGE_MODELS


def _average_quality(results: list[dict[str, object]]) -> float:
    scores: list[float] = []
    for result in results:
        try:
            score = float(result.get("quality_score", 0))
            if score > 1:
                score = score / 100
            scores.append(max(0.0, min(1.0, score)))
        except Exception:
            continue
    if not scores:
        return 0.0
    return sum(scores) / len(scores)


def _collect_failure_modes(results: list[dict[str, object]]) -> list[str]:
    modes: list[str] = []
    for result in results:
        value = result.get("failure_modes", [])
        if isinstance(value, list):
            for item in value:
                text = str(item)
                if text and text not in modes:
                    modes.append(text)
    return modes[:8]


def _collect_next_change(results: list[dict[str, object]]) -> str:
    changes: list[str] = []
    for result in results:
        text = str(result.get("next_pipeline_change", "")).strip()
        if text and text not in changes:
            changes.append(text)
    return " | ".join(changes[:4])


if __name__ == "__main__":
    main()
